"""
Copia el PPTX de diagnóstico TI y sustituye textos por contenido de mesa actuarial,
conservando layout, imágenes y shapes sin text_frame (mismo orden de textos que el dump TI).

Preserva runs en textos para no perder colores de la plantilla; en la lámina
«Veredicto ejecutivo» reaplica tipografía/colores y sincroniza barras (ancho + color).

Reemplazos: plantillas/diagnostico_actuarial_reemplazos.json
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
_WORKSPACE_ROOT = Path(__file__).resolve().parents[2]
REPLACEMENTS_JSON = ROOT / "plantillas" / "diagnostico_actuarial_reemplazos.json"
DEFAULT_SOURCE = (
    _WORKSPACE_ROOT / "Info" / "Diagnostico_TI_La_Internacional_18abr2026_v2.pptx"
)
DEFAULT_OUT = (
    _WORKSPACE_ROOT / "Info" / "Diagnostico_ACT_La_Internacional_18abr2026_v1.pptx"
)

# Diapositiva «Veredicto ejecutivo» (índice 2): pista y barra interna por fila (plantilla TI v2).
_VERDICT_OUTER_BAR_IDXS = (10, 16, 22, 28, 34, 40, 46)
_VERDICT_INNER_BAR_IDXS = (11, 17, 23, 29, 35, 41, 47)
_VERDICT_SCORE_TEXT_IDXS = (12, 18, 24, 30, 36, 42, 48)
_VERDICT_MAIN_SCORE_IDX = 4
_VERDICT_LABEL_IDXS = (9, 15, 21, 27, 33, 39, 45)
_VERDICT_TITLE_IDX = 1
_VERDICT_SUBMAIN_IDXS = (5, 6)
_VERDICT_FOOTER_IDX = 49

_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
_EMERALD = RGBColor(0x10, 0xB9, 0x81)
_SLATE = RGBColor(0x94, 0xA3, 0xB8)
_SLATE200 = RGBColor(0xE2, 0xE8, 0xF0)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_SCORE_GREEN_THRESHOLD = 3.0


def _assign_tf_preserving_runs(tf, new_text: str) -> None:
    """Asigna texto línea a línea conservando el primer run (color/tamaño) de cada párrafo."""
    parts = (new_text or "").split("\n")
    while len(tf.paragraphs) < len(parts):
        tf.add_paragraph()
    for i, part in enumerate(parts):
        p = tf.paragraphs[i]
        if not p.runs:
            p.add_run()
        p.runs[0].text = part
        for j in range(1, len(p.runs)):
            p.runs[j].text = ""
    for j in range(len(parts), len(tf.paragraphs)):
        p = tf.paragraphs[j]
        if p.runs:
            p.runs[0].text = ""
            for k in range(1, len(p.runs)):
                p.runs[k].text = ""
        else:
            p.text = ""


def _text_shapes_in_order(slide):
    return [sh for sh in slide.shapes if sh.has_text_frame]


def _parse_score_cell(s: str) -> float | None:
    t = s.strip().replace(",", ".").replace("TBD", "").strip()
    if not t:
        return None
    try:
        return float(t)
    except ValueError:
        return None


def _sync_verdict_progress_bars(slide) -> None:
    """Ancho y color de barra según nota; promedio en el score central si hace falta."""
    scores: list[float | None] = []
    for idx in _VERDICT_SCORE_TEXT_IDXS:
        sh = slide.shapes[idx]
        if not sh.has_text_frame:
            scores.append(None)
            continue
        scores.append(_parse_score_cell(sh.text_frame.text))

    for outer_i, inner_i, sc in zip(
        _VERDICT_OUTER_BAR_IDXS, _VERDICT_INNER_BAR_IDXS, scores
    ):
        if sc is None:
            continue
        outer = slide.shapes[outer_i]
        inner = slide.shapes[inner_i]
        clamped = max(0.0, min(5.0, sc))
        inner.width = int(outer.width * clamped / 5.0)
        inner.fill.solid()
        inner.fill.fore_color.rgb = (
            _EMERALD if sc >= _SCORE_GREEN_THRESHOLD else _AMBER
        )

    valid = [s for s in scores if s is not None]
    main_sh = slide.shapes[_VERDICT_MAIN_SCORE_IDX]
    if valid and main_sh.has_text_frame:
        m = sum(valid) / len(valid)
        main_txt = _parse_score_cell(main_sh.text_frame.text)
        if main_txt is None or abs(main_txt - m) > 0.05:
            _assign_tf_preserving_runs(main_sh.text_frame, f"{m:.2f}")


def _restyle_verdict_slide(slide) -> None:
    """Fuerza contraste de la lámina 3 según la plantilla TI (título blanco, notas ámbar/verde)."""
    t_tf = slide.shapes[_VERDICT_TITLE_IDX].text_frame
    p0 = t_tf.paragraphs[0]
    if p0.runs:
        r = p0.runs[0]
        r.font.color.rgb = _WHITE
        r.font.bold = True
        r.font.size = Pt(24)

    main_tf = slide.shapes[_VERDICT_MAIN_SCORE_IDX].text_frame
    r = main_tf.paragraphs[0].runs[0]
    r.font.color.rgb = _AMBER
    r.font.bold = True
    r.font.size = Pt(40)

    for idx in _VERDICT_SUBMAIN_IDXS:
        p = slide.shapes[idx].text_frame.paragraphs[0]
        if not p.runs:
            continue
        r = p.runs[0]
        if idx == 5:
            r.font.color.rgb = _SLATE
            r.font.size = Pt(14)
            r.font.bold = False
        else:
            r.font.color.rgb = _AMBER
            r.font.bold = True
            r.font.size = Pt(11)

    for idx in _VERDICT_LABEL_IDXS:
        p = slide.shapes[idx].text_frame.paragraphs[0]
        if not p.runs:
            continue
        r = p.runs[0]
        r.font.color.rgb = _SLATE200
        r.font.size = Pt(12)
        r.font.bold = False

    for idx in _VERDICT_SCORE_TEXT_IDXS:
        tf = slide.shapes[idx].text_frame
        p = tf.paragraphs[0]
        if not p.runs:
            continue
        r = p.runs[0]
        sc = _parse_score_cell(tf.text)
        if sc is not None and sc >= _SCORE_GREEN_THRESHOLD:
            r.font.color.rgb = _EMERALD
        else:
            r.font.color.rgb = _AMBER
        r.font.bold = True
        r.font.size = Pt(13)

    foot = slide.shapes[_VERDICT_FOOTER_IDX].text_frame.paragraphs[0]
    if foot.runs:
        rf = foot.runs[0]
        rf.font.color.rgb = _AMBER
        rf.font.size = Pt(11)
        rf.font.bold = False


def main() -> int:
    source = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_SOURCE
    out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else DEFAULT_OUT
    json_path = Path(sys.argv[3]) if len(sys.argv) > 3 else REPLACEMENTS_JSON

    if not source.is_file():
        print(f"No existe fuente: {source}", file=sys.stderr)
        return 1
    if not json_path.is_file():
        print(f"No existe JSON: {json_path}", file=sys.stderr)
        return 1

    per_slide: list[list[str]] = json.loads(
        json_path.read_text(encoding="utf-8")
    )
    prs = Presentation(str(source))
    if len(prs.slides) != len(per_slide):
        print(
            f"Diapositivas PPTX ({len(prs.slides)}) != slides en JSON ({len(per_slide)})",
            file=sys.stderr,
        )
        return 1

    for si, slide in enumerate(prs.slides):
        shapes_tf = _text_shapes_in_order(slide)
        texts = per_slide[si]
        if len(shapes_tf) != len(texts):
            print(
                f"Slide {si}: shapes con texto ({len(shapes_tf)}) != cadenas JSON ({len(texts)})",
                file=sys.stderr,
            )
            return 1
        for sh, t in zip(shapes_tf, texts):
            _assign_tf_preserving_runs(sh.text_frame, t)
        if si == 2:
            _sync_verdict_progress_bars(slide)
            _restyle_verdict_slide(slide)

    prs.save(str(out_path))
    print(out_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
