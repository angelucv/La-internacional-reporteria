"""
Genera un .pptx por mesa a partir de los borradores Markdown en
reuniones/2026/reportes-asesoria/REP-POR-*-Semana01-2026-BORRADOR.md

Colores alineados al criterio corporativo (#1B3A5C títulos).
"""
from __future__ import annotations

import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
_WORKSPACE_ROOT = Path(__file__).resolve().parents[2]
REPORT_DIR = ROOT / "reuniones" / "2026" / "reportes-asesoria"
GLOB_MD = "REP-POR-*-Semana01-2026-BORRADOR.md"
# Copia opcional de la presentación actuarial Semana 1 a carpeta Semanal (Mesa Actuarial)
ACT_PPTX_COPY_DIR = _WORKSPACE_ROOT / "Reportes" / "Mesa Actuarial" / "Semanal"

TITLE_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
BODY_GRAY = RGBColor(0x33, 0x33, 0x33)


def _strip_md(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    return s.strip()


def _parse_sections(md: str) -> dict[str, str]:
    sections: dict[str, list[str]] = {}
    current: str | None = None
    skip_until_h2 = True
    for line in md.splitlines():
        if line.startswith("## "):
            skip_until_h2 = False
            current = line[3:].strip()
            sections[current] = []
            continue
        if skip_until_h2 or current is None:
            continue
        if line.startswith("# ") and not line.startswith("##"):
            continue
        sections[current].append(line)
    return {k: "\n".join(v).strip() for k, v in sections.items()}


def _meta_line(md: str, label: str) -> str:
    for line in md.splitlines():
        if label in line and "|" in line:
            parts = [p.strip() for p in line.split("|")]
            parts = [p for p in parts if p]
            if len(parts) >= 2 and label in parts[0]:
                return _strip_md(parts[1])
    return ""


def _parse_hallazgos_table(block: str) -> list[tuple[str, str, str]]:
    rows: list[tuple[str, str, str]] = []
    for line in block.splitlines():
        line = line.strip()
        if (
            not line.startswith("|")
            or line.startswith("|:---")
            or ("Tema" in line and "Hallazgo" in line)
        ):
            continue
        cells = [c.strip() for c in line.split("|")]
        cells = [c for c in cells if c]
        if len(cells) >= 3:
            rows.append((_strip_md(cells[0]), _strip_md(cells[1]), _strip_md(cells[2])))
        elif len(cells) == 2:
            rows.append((_strip_md(cells[0]), _strip_md(cells[1]), ""))
    return rows


def _bullets_from_block(block: str) -> list[str]:
    items: list[str] = []
    for line in block.splitlines():
        line = line.strip()
        if line.startswith("- "):
            items.append(_strip_md(line[2:]))
        elif line.startswith("*") and not line.startswith("**"):
            items.append(_strip_md(line.lstrip("* ").strip()))
    return items


def _set_title(tf, text: str, *, size: int = 28) -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TITLE_BLUE
    p.alignment = PP_ALIGN.LEFT


def _add_body_slide(prs: Presentation, title: str, paragraphs: list[str]) -> None:
    layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(layout)
    _set_title(slide.shapes.title.text_frame, title, size=22)
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, para in enumerate(paragraphs):
        if not para:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        p.font.size = Pt(14)
        p.font.color.rgb = BODY_GRAY
        p.level = 0
        p.space_after = Pt(6)


def _add_table_slide(prs: Presentation, title: str, rows: list[tuple[str, str, str]]) -> None:
    try:
        blank = prs.slide_layouts[6]
    except IndexError:
        blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    _set_title(box.text_frame, title, size=22)
    if not rows:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(1))
        tb.text_frame.paragraphs[0].text = "(Sin filas en tabla — completar al recibir minutas e insumos.)"
        return
    nrows = len(rows) + 1
    ncols = 3
    left, top = Inches(0.5), Inches(1.2)
    width = Inches(9)
    height = Inches(0.35 * min(nrows + 1, 12))
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    hdr = ("Tema", "Hallazgo (borrador)", "Riesgo")
    for c, h in enumerate(hdr):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            if c >= ncols:
                break
            cell = table.cell(r, c)
            cell.text = val[:500] if val else "—"
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)


def build_pptx_from_md(md_path: Path, out_path: Path) -> None:
    text = md_path.read_text(encoding="utf-8")
    mesa = _meta_line(text, "**Mesa**") or "Mesa"
    ref = _meta_line(text, "**Referencia interna**") or ""
    oleada = _meta_line(text, "**Oleada**") or ""
    responsables = _meta_line(text, "**Responsables (Solicitud)**") or ""

    sections = _parse_sections(text)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Evaluación de gestión, hallazgos y controles"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_BLUE
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(26)
    st = slide.placeholders[1].text_frame
    st.clear()
    sub = f"{mesa}\nPlan de Optimización y Rentabilidad 2026\n{oleada}\n{ref}\n\nClasificación: Confidencial"
    p = st.paragraphs[0]
    p.text = sub
    p.font.size = Pt(14)
    p.font.color.rgb = BODY_GRAY
    if responsables:
        p2 = st.add_paragraph()
        p2.text = f"Responsables (Solicitud): {responsables}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = BODY_GRAY

    # Objetivo + aspectos
    s1 = sections.get("1. Objetivo (primera semana)", "")
    s2 = sections.get("2. Aspectos generales y reunión inicial", "")
    if s1 or s2:
        body_parts = []
        if s1:
            body_parts.append(_strip_md(s1.replace("\n", " ")))
        if s2:
            body_parts.extend(_bullets_from_block(s2) or [_strip_md(s2.replace("\n", " "))])
        _add_body_slide(prs, "Objetivo y aspectos generales", body_parts)

    # Hallazgos
    s3 = sections.get("3. Hallazgos (borrador)", "")
    intro, _, table_part = s3.partition("| Tema |")
    rows = _parse_hallazgos_table("| Tema |" + table_part if table_part else s3)
    if intro.strip() and not intro.strip().startswith("|"):
        intro_clean = _strip_md(intro.replace("\n", " "))
        if len(intro_clean) > 30:
            _add_body_slide(prs, "Hallazgos — contexto", [intro_clean[:1500]])
    _add_table_slide(prs, "Hallazgos (borrador)", rows)

    # Sección 4 (perspectiva / acciones — cualquier título "4. …")
    s4 = ""
    for k in sections:
        if k.startswith("4."):
            s4 = sections[k]
            break
    bullets4 = _bullets_from_block(s4) or ([_strip_md(s4)] if s4 else [])
    s5 = sections.get("5. Requerimientos Solicitud v5 — recorte Semana 1 (mesa financiera)", "")
    s5b = sections.get("5. Requerimientos Solicitud v5 — recorte Semana 1 (mesa actuarial)", "")
    s5c = sections.get("5. Requerimientos Solicitud v5 — recorte Semana 1 (enfoque legal)", "")
    s5d = sections.get("5. Requerimientos Solicitud v5 — recorte Semana 1 (cumplimiento)", "")
    s5e = sections.get("5. Requerimientos Solicitud v5 — enlace Semana 1 (cronograma agregado)", "")
    s5f = sections.get("5. Requerimientos Solicitud v5 — recorte Semana 1 / 1–2 (tecnología)", "")
    s5_any = s5 or s5b or s5c or s5d or s5e or s5f
    for k in sections:
        if k.startswith("5.") and not s5_any:
            s5_any = sections[k]
            break
    bullets5 = _bullets_from_block(s5_any) if s5_any else []
    if s5_any and not bullets5:
        para = _strip_md(s5_any.replace("\n", " "))
        if para:
            bullets5 = [para]

    acc_paras = []
    if bullets4:
        acc_paras.append("Perspectiva (siguiente oleada):")
        acc_paras.extend(f"• {b}" for b in bullets4)
    if bullets5:
        acc_paras.append("Ámbitos de la solicitud (recorte):")
        acc_paras.extend(f"• {b}" for b in bullets5)
    if acc_paras:
        _add_body_slide(prs, "Perspectiva y marco documental", acc_paras)

    # Sección opcional 6 (p. ej. control semanal — título = texto tras "6.")
    s6_keys = [k for k in sections if k.startswith("6.")]
    for key in sorted(s6_keys):
        s6 = sections.get(key, "")
        parts6 = _bullets_from_block(s6) or ([_strip_md(s6.replace("\n", " "))] if s6 else [])
        if parts6:
            title6 = key.split(".", 1)[-1].strip() if "." in key else key.strip()
            _add_body_slide(prs, title6, parts6[:24])

    prs.save(out_path)


def main() -> None:
    paths = sorted(REPORT_DIR.glob(GLOB_MD))
    if not paths:
        print(f"No se encontraron archivos {GLOB_MD} en {REPORT_DIR}", file=sys.stderr)
        sys.exit(1)
    for md in paths:
        out = md.with_suffix(".pptx")
        build_pptx_from_md(md, out)
        print(f"Generado: {out}")

    act_pptx = REPORT_DIR / "REP-POR-ACT-Semana01-2026-BORRADOR.pptx"
    act_md = REPORT_DIR / "REP-POR-ACT-Semana01-2026-BORRADOR.md"
    if ACT_PPTX_COPY_DIR.parent.is_dir():
        ACT_PPTX_COPY_DIR.mkdir(parents=True, exist_ok=True)
        if act_pptx.is_file():
            dest = ACT_PPTX_COPY_DIR / act_pptx.name
            shutil.copy2(act_pptx, dest)
            print(f"Copiado a Semanal: {dest}")
        if act_md.is_file():
            shutil.copy2(act_md, ACT_PPTX_COPY_DIR / act_md.name)
            print(f"Copiado a Semanal: {ACT_PPTX_COPY_DIR / act_md.name}")
    else:
        print(
            f"Aviso: no existe la ruta padre de Semanal, no se copió: {ACT_PPTX_COPY_DIR}",
            file=sys.stderr,
        )


if __name__ == "__main__":
    main()
