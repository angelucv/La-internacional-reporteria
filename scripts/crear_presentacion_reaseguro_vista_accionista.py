"""Genera PPTX ejecutivo — reaseguro para nuevo accionista (no operativo).



Salida:

  Reportes/Mesa Actuarial/Materiales/La_Internacional_Reaseguro_Vista_Accionista_2026.pptx



Contenido alineado a:

  Materiales/MATRIZ-REASEGURO-EJECUTIVA-MONTOS-CONTRATOS-Y-PA-2026-04-23.md

"""



from __future__ import annotations



from pathlib import Path



from pptx import Presentation

from pptx.util import Pt





ROOT_REPORTERIA = Path(__file__).resolve().parents[1]

WORKSPACE = ROOT_REPORTERIA.parent

OUT_DIR = WORKSPACE / "Reportes" / "Mesa Actuarial" / "Materiales"

OUT_PPTX = OUT_DIR / "La_Internacional_Reaseguro_Vista_Accionista_2026.pptx"





def _bullets_slide(prs: Presentation, title: str, lines: list[str], *, size: int = 14) -> None:

    layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame

    body.clear()

    for i, line in enumerate(lines):

        p = body.paragraphs[0] if i == 0 else body.add_paragraph()

        p.text = line

        p.level = 0

        p.font.name = "Arial"

        p.font.size = Pt(size)





def build() -> None:

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    prs.slide_width = int(prs.slide_width)

    prs.slide_height = int(prs.slide_height)



    # Portada (layout título)

    slide0 = prs.slides.add_slide(prs.slide_layouts[0])

    slide0.shapes.title.text = "Reaseguro — vista para accionista"

    sub = slide0.placeholders[1]

    sub.text = (

        "La Internacional de Seguros, S.A.\n"

        "Síntesis de concentración, cifras 2025 y trámite regulatorio\n"

        "Abril 2026 · Confidencial"

    )

    for p in sub.text_frame.paragraphs:

        p.font.name = "Arial"

        p.font.size = Pt(18)



    _bullets_slide(

        prs,

        "Qué debe retener la junta",

        [

            "El reaseguro distribuye el riesgo técnico y fija el gasto en prima cedida; en 2025 el peso en dólares "

            "está en tratados no proporcionales y facultativos de salud (Redbridge) además del programa Bouquet.",

            "En bolívares el facultativo concentra casi la totalidad del gasto de reaseguro del año según el resumen "

            "interno de pagos — conviene validar conciliación con contabilidad.",

            "La gobernanza relevante es la alineación entre contratos PDF (Legal), pagos (Finanzas) y la relación "

            "declarada a SUDEASEG (fechas de remisión como huella administrativa).",

        ],

        size=13,

    )



    _bullets_slide(

        prs,

        "Acumulación de prima neta pagada (2025)",

        [

            "USD — Bouquet proporcional: 146.993 (Bee, Delta, Hannover/Kairos, Kairos, R.I.V).",

            "USD — No proporcional salud (Redbridge): 4,36 millones.",

            "USD — Facultativo puro: 2,98 millones (incluye Redbridge y otras contrapartes).",

            "USD — Total informado en resumen: 7,48 millones.",

            "Bs — Bouquet proporcional: 121.412 (Delta, Kairos).",

            "Bs — Facultativo: 27,33 millones (una línea en fuente; revisar etiqueta con Finanzas).",

            "Fuente: Reportes/Resumen de pagos Reaseguros año 2025 Juan.xlsx.",

        ],

        size=12,

    )



    _bullets_slide(

        prs,

        "Contratos y declaración al regulador",

        [

            "Contratos PDF: Mesa Actuarial/Materiales/De-Gerencia-Legal/01-Reaseguro-Contratos "

            "(familia Bouquet 2026, proporcionales auto/fianzas/HCM, vida AON).",

            "Relación SUDEASEG 2025: Carpeta compartida Actuarial — archivo "

            "«f- Relación de Contratos de Reaseguro Año 2025 SUDEASEG.xlsx» (hojas 2025 Bs y USD).",

            "Matriz detallada con tablas y riesgos: MATRIZ-REASEGURO-EJECUTIVA-MONTOS-CONTRATOS-Y-PA-2026-04-23.md.",

        ],

        size=12,

    )



    _bullets_slide(

        prs,

        "Procedimiento administrativo (marco)",

        [

            "Remisión de la relación de contratos de reaseguro al regulador; en el Excel constan fechas "

            "(ej.: 31-mar-2025, 18-jun-2025, 22-oct-2025) según líneas y endosos.",

            "Expediente complementario en De-Gerencia-Legal/06-Sancionatorios-SUDEASEG para comunicaciones "

            "y procedimientos supervisoras — coordinar lectura con Legal antes de representaciones públicas.",

            "Cruce Legal ↔ SUDEASEG documentado en MATRIZ-REASEGURO-Legal-vs-SUDEASEG-2026-04-24.md.",

        ],

        size=12,

    )



    prs.save(str(OUT_PPTX))

    print(OUT_PPTX)





if __name__ == "__main__":

    build()


