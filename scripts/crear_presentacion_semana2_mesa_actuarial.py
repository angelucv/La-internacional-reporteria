"""

A partir del borrador Semana 1, genera versiones Semana 2 (Mesa Actuarial):

- La_Internacional_Semana2_2026.pptx

- REP-POR-ACT-Semana02-2026-BORRADOR.pptx



Portada: semana 2, referencia Semana02, fechas.

Diapositiva 5: el texto de avance se escribe en el **placeholder de cuerpo** (no en un textbox

superpuesto), porque el layout original tapaba el cuadro añadido y no se veía en PowerPoint.



Diapositivas 2 a 4: se **reemplazan** por contenido **exclusivo de Semana 2** (objetivo y hechos,

tabla de hallazgos con columnas **Riesgo | Ámbito | Hallazgo** al estilo de la lámina Mesa Actuarial

de `La_Internacional_Semana1_2026`, perspectiva y cierre de la semana). No arrastran texto de la

Semana 1.



Diapositivas nuevas al final:

- Detalle del vaciado de Gerencia Legal en Mesa Actuarial.

- Detalle actuarial (promoción, minuta, matriz, acuerdos).

"""

from __future__ import annotations



import re

import shutil

from pathlib import Path



from pptx import Presentation

from pptx.util import Pt



ROOT_REPORTERIA = Path(__file__).resolve().parents[1]

WORKSPACE = ROOT_REPORTERIA.parent

SEMANAL = WORKSPACE / "Reportes" / "Mesa Actuarial" / "Semanal"

SRC = SEMANAL / "REP-POR-ACT-Semana01-2026-BORRADOR.pptx"

OUT_LA = SEMANAL / "La_Internacional_Semana2_2026.pptx"

OUT_REP = SEMANAL / "REP-POR-ACT-Semana02-2026-BORRADOR.pptx"





def _replace_in_slide_text(slide, mapping: list[tuple[str, str]]) -> None:

    for shape in slide.shapes:

        if not shape.has_text_frame:

            continue

        for para in shape.text_frame.paragraphs:

            t = para.text

            t = re.sub(

                r"Semana 1 \([^)]+\)",

                "Semana 2 (21–25 abr. 2026)",

                t,

            )

            for old, new in mapping:

                t = t.replace(old, new)

            para.text = t





def _text_frame_set_paragraphs(tf, lines: list[str], *, body_size: int = 11) -> None:

    """Reemplaza el contenido del TextFrame por una lista de párrafos (nivel 0)."""

    while len(tf.paragraphs) > 1:

        el = tf.paragraphs[-1]._element

        el.getparent().remove(el)

    if not lines:

        tf.paragraphs[0].text = ""

        return

    tf.paragraphs[0].text = lines[0]

    for para in tf.paragraphs:

        para.level = 0

        para.font.name = "Arial"

        para.font.size = Pt(body_size)

        for run in para.runs:

            run.font.name = "Arial"

            run.font.size = Pt(body_size)

    for line in lines[1:]:

        p = tf.add_paragraph()

        p.text = line

        p.level = 0

        p.font.name = "Arial"

        p.font.size = Pt(body_size)

        for run in p.runs:

            run.font.name = "Arial"

            run.font.size = Pt(body_size)





def _placeholder_body(slide, idx: int = 1):
    """Devuelve la forma placeholder de cuerpo (típ. idx 1) o None."""
    for shape in slide.shapes:
        try:
            if shape.is_placeholder and shape.placeholder_format.idx == idx:
                if shape.has_text_frame:
                    return shape
        except ValueError:
            continue
    return None


def _week2_fill_slide_objective_and_context(slide) -> None:
    """Diapositiva 2: objetivo y hechos de la semana 2 únicamente (no arrastra Semana 1)."""
    body = _placeholder_body(slide, 1)
    if body is None:
        return
    lines = [
        "Objetivo (Semana 2): consolidar la lectura integrada del pliego frente a lo ya depositado "
        "y promovido en la carpeta Mesa Actuarial: traza entre la bandeja de gerencia actuarial y el "
        "paquete para Drive / agente compilador; focalizar vacíos que aún impiden un juicio técnico "
        "completo (SEFA, primas transaccionales, reaseguro contractual, homologación SUDEASEG en siniestros).",
        "",
        "Aspectos generales (solo esta semana):",
        "• Promoción masiva desde Reportes/Carpeta compartida Actuarial hacia "
        "Mesa Actuarial/Materiales/Promocion-2026-04-23-Carpeta-compartida-Actuarial (176 archivos), "
        "con INDICE-PROMOCION-CARPETA-ACTUARIAL-A-MESA.md y actualización del documento maestro del workspace.",
        "• Publicación de la minuta de validación vs Solicitud v2 (MIN-POR-ACT-20260424-04; Word con fecha "
        "de actualización en el nombre en Minutas Reuniones/).",
        "• Presentación de control Semana 2 a 7 diapositivas: síntesis visible en el cuerpo de la lámina 5; "
        "láminas 6–7 con vaciado Legal y detalle Actuarial.",
        "• Integración del árbol de Reportes/GERENCIA LEGAL en Mesa Actuarial (176 archivos en espejo íntegro "
        "y clasificación De-Gerencia-Legal/01–13 con INDICE-VACIADO-Y-MAPA-A-MESA-ACTUARIAL.md), sin alterar el origen.",
    ]
    _text_frame_set_paragraphs(body.text_frame, lines, body_size=11)


def _week2_fill_hallazgos_table(slide) -> None:
    """Diapositiva 3: tabla al estilo Mesa Actuarial de Semana 1 (riesgo visible primero), solo Semana 2."""
    tbl = None
    for shape in slide.shapes:
        if hasattr(shape, "table"):
            tbl = shape.table
            break
    if tbl is None:
        return
    # Misma cantidad de filas que la plantilla Semana 1 (encabezado + 8 filas de datos).
    data = [
        ("Riesgo", "Ámbito (Mesa Actuarial — Sem. 2)", "Hallazgo (solo hechos de la semana)"),
        (
            "Bajo",
            "Bandeja actuarial en Mesa",
            "176 archivos promovidos desde Carpeta compartida Actuarial; índice "
            "INDICE-PROMOCION-CARPETA-ACTUARIAL-A-MESA.md; documento maestro del workspace actualizado.",
        ),
        (
            "Bajo",
            "Minuta vs Solicitud v2",
            "Código MIN-POR-ACT-20260424-04; archivo Word con fecha de actualización en el nombre "
            "(2026-04-24) en Minutas Reuniones/.",
        ),
        (
            "Bajo",
            "Presentación de control semanal",
            "La_Internacional_Semana2_2026 y REP-POR-ACT-Semana02-2026-BORRADOR: 7 diapositivas; "
            "síntesis en placeholder de cuerpo (lámina 5); láminas 6–7 ampliación Legal y Actuarial.",
        ),
        (
            "Bajo",
            "Gerencia Legal en la mesa",
            "Espejo íntegro vaciado-completo-2026-04-23 (176 archivos); copias temáticas "
            "De-Gerencia-Legal/01–13; mapa INDICE-VACIADO-Y-MAPA-A-MESA-ACTUARIAL.md.",
        ),
        (
            "Bajo",
            "Certificación y reportería de cierre",
            "Reservas: PDF y series en carpeta compartida actuarial (certificación, ONR, reserva "
            "complementaria). Cierres: consolidación manual hacia Reportes/Mesa con control humano "
            "explícito y aceptable en la semana.",
        ),
        (
            "Alto",
            "Automatización reportería / BI",
            "No hay desarrollo ni tubería que publique cierres de gestión de forma recurrente hacia "
            "Mesa Actuarial ni hacia los tableros BI.",
        ),
        (
            "Medio",
            "Regulador, analítica y reaseguro",
            "Cola productos SUDEASEG en espera; condicionados en De-Gerencia-Legal/02; recibos de "
            "prima cubiertos; inventario pendiente de otros formatos analíticos con Tecnología/Finanzas; "
            "SEFA y primas transaccionales como vacíos del pliego. PDF de reaseguro en Mesa + matriz "
            "Legal vs SUDEASEG 2025; cruce sistemático en curso.",
        ),
        (
            "Bajo",
            "Drive / compilador POR",
            "Carpeta Mesa Actuarial más homogénea; minuta y matriz orientan al compilador multi-mesa "
            "y al paquete citado para Drive.",
        ),
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            tbl.rows[r].cells[c].text = val


def _week2_fill_perspective_slide(slide) -> None:
    """Diapositiva 4: perspectiva y cierre de la semana 2 (sin texto de «siguiente oleada» de Semana 1)."""
    body = _placeholder_body(slide, 1)
    if body is None:
        return
    lines = [
        "Perspectiva inmediata (Semana 2):",
        "• Profundizar densidad transaccional y evidencia analítica (SEFA, primas, cuadres al regulador).",
        "• Sostener conclusiones con el corpus contractual de reaseguro ya ubicado en Mesa.",
        "",
        "Control semanal — balance (solo Semana 2):",
        "• Positivo en integración documental: promoción actuarial + vaciado Legal + minuta y matriz de reaseguro "
        "alineadas a la Solicitud v2.",
        "• Mixto en cierre de brechas del pliego hasta coordinación con Finanzas y Tecnología.",
        "• Mirada hacia adelante: paquete citado/versionado para Drive; hilo para compilador multi-mesa del POR; "
        "cruce sistemático contratos Legal vs tabla SUDEASEG.",
    ]
    _text_frame_set_paragraphs(body.text_frame, lines, body_size=11)


def _slide5_fill_control_placeholder(slide) -> None:

    """Escribe la síntesis Semana 2 en el placeholder de cuerpo (idx 1), visible en PowerPoint."""

    disclaimer = (

        "Este borrador sirve de control semanal informativo de la posición de la mesa; puede "

        "complementarse con el registro formal de sesiones y el expediente de la mesa sin sustituirlos.*"

    )

    body = None

    for shape in slide.shapes:

        try:

            if shape.is_placeholder and shape.placeholder_format.idx == 1:

                body = shape

                break

        except ValueError:

            continue

    if body is None or not body.has_text_frame:

        return

    lines = [

        "Síntesis Semana 2 — Mesa Actuarial (ampliación en las dos últimas diapositivas de esta presentación).",

        "• Insumos gerencia actuarial: 176 archivos en Mesa Actuarial/Materiales/Promocion-2026-04-23-Carpeta-compartida-Actuarial; traza INDICE-PROMOCION-CARPETA-ACTUARIAL-A-MESA.md.",

        "• Gerencia Legal integrada: espejo en Mesa Actuarial/Gerencia-Legal/vaciado-completo-2026-04-23 (176 archivos) y clasificación por tema en Materiales/De-Gerencia-Legal/ (carpetas 01 a 13); mapa INDICE-VACIADO-Y-MAPA-A-MESA-ACTUARIAL.md.",

        "• Minuta vs Solicitud v2: código MIN-POR-ACT-20260424-04; archivo MIN-POR-ACT-20260424-Mesa-Actuarial-vs-Solicitud_actualizado-2026-04-24.docx en Minutas Reuniones/.",

        "• Matriz reaseguro (contratos Legal vs relación SUDEASEG 2025): Mesa Actuarial/Materiales/MATRIZ-REASEGURO-Legal-vs-SUDEASEG-2026-04-24.md.",

        "• Acuerdos de foco: inventario carga analítica distinta de recibos; SEFA y primas con Finanzas/Tecnología; cola productos SUDEASEG en espera; paquete para Drive y agente compilador POR.",

        "",

        disclaimer,

    ]

    _text_frame_set_paragraphs(body.text_frame, lines, body_size=11)

    # Disclaimer en cursiva un poco más pequeña

    last = body.text_frame.paragraphs[-1]

    last.font.size = Pt(10)

    last.font.italic = True

    for run in last.runs:

        run.font.italic = True

        run.font.size = Pt(10)

        run.font.name = "Arial"





def _append_title_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:

    layout = prs.slide_layouts[1]  # Title and Content

    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = title

    body = slide.placeholders[1]

    _text_frame_set_paragraphs(body.text_frame, bullets, body_size=12)





def _append_legal_detail_slides(prs: Presentation) -> None:

    _append_title_content_slide(

        prs,

        "Gerencia Legal — vaciado en Mesa Actuarial (2026-04-23)",

        [

            "Origen en disco: Reportes/GERENCIA LEGAL (sin modificar). Todo el árbol se replicó en la carpeta de trabajo de la mesa.",

            "Espejo íntegro (misma jerarquía que Legal):",

            "Mesa Actuarial/Gerencia-Legal/vaciado-completo-2026-04-23/",

            "Clasificación por tema para la mesa (copias físicas bajo Materiales/):",

            "Mesa Actuarial/Materiales/De-Gerencia-Legal/",

            "01-Reaseguro-Contratos | 02-Productos-Condicionados-SUDEASEG | 03-Intermediarios-Sucursales | 04-Canales-Alternativos",

            "05-Plan-Incentivo-esquema-comercial | 06-Sancionatorios-SUDEASEG | 07a-Cumplimiento-AML-FT | 07b-Juicios-activos",

            "08-Contribuciones-aportes | 09-Solvencias-parafiscales | 10-Presupuesto-y-gobierno-corporativo",

            "11a-Arrendamientos | 11b-Contratos-proveedores | 12-Declaraciones-ISR | 13-Propiedad-intelectual",

            "Índice y mapa de correspondencia: Mesa Actuarial/Gerencia-Legal/INDICE-VACIADO-Y-MAPA-A-MESA-ACTUARIAL.md",

            "Cruce reaseguro: use la carpeta 01-Reaseguro-Contratos junto con MATRIZ-REASEGURO-Legal-vs-SUDEASEG-2026-04-24.md y el Excel f- Relación de Contratos de Reaseguro Año 2025 SUDEASEG en Carpeta compartida Actuarial.",

        ],

    )





def _append_actuarial_detail_slide(prs: Presentation) -> None:

    _append_title_content_slide(

        prs,

        "Actuarial — minuta, hallazgos y compilación POR",

        [

            "Promoción desde Carpeta compartida Actuarial: 176 archivos en Materiales/Promocion-2026-04-23-Carpeta-compartida-Actuarial.",

            "Minuta de validación vs Solicitud consolidada Definitivo v2: MIN-POR-ACT-20260424-04; documento Word con fecha de actualización en el nombre (2026-04-24).",

            "Hallazgos relevantes en matriz: reportería de gestión manual (riesgo Bajo); automatización de reportería/BI hacia Mesa Actuarial (Alto — sin desarrollo a la fecha).",

            "Carga analítica al regulador: cubiertos recibos de prima (100RPCD); pendiente inventario de otros formatos analíticos (no recibos) con Tecnología/Finanzas.",

            "Oferta / cola de productos SUDEASEG: en espera de confirmación; condicionados disponibles en De-Gerencia-Legal/02-Productos-Condicionados-SUDEASEG.",

            "Próximo paso: sincronizar Mesa Actuarial con Drive y alimentar al agente compilador multi-mesa del POR con nombres de archivo y referencias citadas en índices.",

        ],

    )





def build() -> None:

    if not SRC.is_file():

        raise FileNotFoundError(SRC)

    tmp = SEMANAL / "_tmp_semana2_build.pptx"

    shutil.copy2(SRC, tmp)

    prs = Presentation(str(tmp))



    mapping = [

        ("REP-POR-ACT-Semana01-2026-BORRADOR", "REP-POR-ACT-Semana02-2026-BORRADOR"),

        ("Semana01", "Semana02"),

    ]

    for slide in prs.slides:

        _replace_in_slide_text(slide, mapping)



    if len(prs.slides) >= 2:

        prs.slides[1].shapes.title.text = "Objetivo y aspectos generales (Semana 2)"

        _week2_fill_slide_objective_and_context(prs.slides[1])



    if len(prs.slides) >= 3:

        s3 = prs.slides[2]

        if s3.shapes.title is not None:

            s3.shapes.title.text = "Hallazgos Mesa Actuarial (Semana 2)"

        elif len(s3.shapes) > 0 and s3.shapes[0].has_text_frame:

            s3.shapes[0].text_frame.paragraphs[0].text = "Hallazgos Mesa Actuarial (Semana 2)"

        _week2_fill_hallazgos_table(s3)



    if len(prs.slides) >= 4:

        prs.slides[3].shapes.title.text = "Perspectiva y control (Semana 2)"

        _week2_fill_perspective_slide(prs.slides[3])



    if len(prs.slides) >= 5:

        _slide5_fill_control_placeholder(prs.slides[4])



    _append_legal_detail_slides(prs)

    _append_actuarial_detail_slide(prs)



    for out in (OUT_LA, OUT_REP):

        prs.save(str(out))

    tmp.unlink(missing_ok=True)

    print(OUT_LA)

    print(OUT_REP)

    print(f"Total diapositivas: {len(prs.slides)}")





if __name__ == "__main__":

    build()


